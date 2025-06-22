import io
import os
from typing import Dict, Any, List

import pptx
from pptx.exc import PackageNotFoundError
import fitz  # PyMuPDF
from google.cloud import storage

from ..state_models import SlideContent, DocumentAnalysisResult

def _extract_text_from_pptx(blob: bytes) -> List[SlideContent]:
    """.pptxファイルからスライドごとのテキストを抽出する"""
    slides_content: List[SlideContent] = []
    try:
        prs = pptx.Presentation(io.BytesIO(blob))
        for i, slide in enumerate(prs.slides):
            text_runs = []
            # シェイプからテキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            # ノートからテキストを抽出
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide else ""
            
            slides_content.append(SlideContent(
                slide_number=i + 1,
                text='\n'.join(text_runs),
                notes=notes_text,
                title=slide.shapes.title.text if slide.has_notes_slide and slide.shapes.title else None
            ))
    except PackageNotFoundError:
        raise ValueError("無効なPowerPointファイル形式です。")
    return slides_content

def _extract_text_from_pdf(blob: bytes) -> List[SlideContent]:
    """PDFファイルからページごとのテキストを抽出する"""
    slides_content: List[SlideContent] = []
    with fitz.open(stream=blob, filetype="pdf") as doc:
        for i, page in enumerate(doc):
            slides_content.append(SlideContent(
                slide_number=i + 1,
                text=page.get_text("text")
            ))
    return slides_content

def _download_from_gcs(gcs_path: str) -> bytes:
    """GCSからファイルをダウンロードしてbytesとして返す"""
    try:
        client = storage.Client()
        # gcs_path is expected to be in "gs://bucket-name/path/to/object" format
        bucket_name, blob_name = gcs_path.replace("gs://", "").split("/", 1)
        bucket = client.bucket(bucket_name)
        blob = bucket.blob(blob_name)
        return blob.download_as_bytes()
    except Exception as e:
        # NOTE: In a real application, consider more specific exception handling
        # and ensure the Cloud Run service account has "Storage Object Viewer" role.
        raise RuntimeError(f"GCSからのファイルダウンロードに失敗しました: {gcs_path}. エラー: {e}")

def parse_presentation_document(file_path: str) -> Dict[str, Any]:
    """
    GCS上のプレゼンテーションファイル(.pptx, .pdf)を解析し、
    テキストコンテンツをJSON形式で返す。

    Args:
        file_path: GCS上のファイルパス (例: "gs://bucket/file.pptx")

    Returns:
        DocumentAnalysisResultモデルに対応する辞書。
    """
    file_name = os.path.basename(file_path)
    try:
        file_blob = _download_from_gcs(file_path)
        
        if file_path.lower().endswith(".pptx"):
            slides = _extract_text_from_pptx(file_blob)
        elif file_path.lower().endswith(".pdf"):
            slides = _extract_text_from_pdf(file_blob)
        else:
            return DocumentAnalysisResult(
                file_name=file_name,
                total_slides=0,
                slides=[],
                error="サポートされていないファイル形式です。.pptxまたは.pdfをアップロードしてください。"
            ).model_dump()

        return DocumentAnalysisResult(
            file_name=file_name,
            total_slides=len(slides),
            slides=slides,
        ).model_dump()

    except Exception as e:
        return DocumentAnalysisResult(
            file_name=file_name,
            total_slides=0,
            slides=[],
            error=f"ファイルの解析中にエラーが発生しました: {str(e)}"
        ).model_dump()